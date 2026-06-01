"""
Generate ART SER Corporate Presentation — High-end, dark-themed, advanced layout.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os, json

# ── Paths ──
BASE = os.path.dirname(os.path.abspath(__file__))
DATA = os.path.join(BASE, "data")
PUBLIC = os.path.join(BASE, "public")
OUT = os.path.join(BASE, "ART_SER_Corporate_Presentation.pptx")

def load(name):
    with open(os.path.join(DATA, name), "r", encoding="utf-8") as f:
        return json.load(f)

company = load("company.json")
clients = load("clients.json")["items"]
services = load("services.json")["items"]
projects = load("projects.json")["items"]
certs = load("certifications.json")["items"]
suppliers = load("suppliers.json")["items"]

# ── Theme colors ──
BG_DARK     = RGBColor(0x0C, 0x0C, 0x0E)
BG_SURFACE  = RGBColor(0x14, 0x14, 0x18)
BG_CARD     = RGBColor(0x1C, 0x1C, 0x22)
ACCENT      = RGBColor(0xC8, 0x96, 0x2E)  # gold
ACCENT_DIM  = RGBColor(0xA0, 0x78, 0x22)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
MUTED       = RGBColor(0x88, 0x88, 0x88)
DIVIDER     = RGBColor(0x2A, 0x2A, 0x30)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, line_color=None, corner_radius=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
                                  left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    return shp


def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_line(slide, left, top, width, color=ACCENT, thickness=2):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(thickness))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def add_accent_bar(slide, left, top, height=Inches(0.6)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Pt(4), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    return bar


def safe_image(path):
    """Return absolute path if image exists, else None."""
    full = os.path.join(PUBLIC, path.lstrip("/"))
    return full if os.path.isfile(full) else None


# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE / COVER
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# Top accent line
add_line(slide, Inches(0), Inches(0), W, ACCENT, 4)

# Company name
add_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
         "ART SER", size=60, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Segoe UI Light")

# Tagline
add_text(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.6),
         company["tagline"]["en"], size=22, color=ACCENT, bold=False, alignment=PP_ALIGN.CENTER)

# Divider
add_line(slide, Inches(5.5), Inches(3.8), Inches(2.3), ACCENT, 2)

# Subtitle
add_text(slide, Inches(2), Inches(4.2), Inches(9), Inches(1),
         "Corporate Capabilities & Client Partnership Overview", size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Legal name + location
add_text(slide, Inches(2), Inches(5.8), Inches(9), Inches(0.4),
         "ART SER DI SHEHEZAD TARIQ  ·  Verona, Italy  ·  Est. 2015  ·  26+ Years Industry Experience",
         size=11, color=MUTED, alignment=PP_ALIGN.CENTER)

# Bottom accent line
add_line(slide, Inches(0), H - Pt(4), W, ACCENT, 4)


# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — OUR CLIENTS (logos + details)
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, Inches(0), Inches(0), W, ACCENT, 3)

# Section header
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
         "OUR CLIENTS", size=11, color=ACCENT, bold=True, font_name="Segoe UI Semibold")
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
         "Trusted Partnerships with Industry Leaders", size=32, color=WHITE, bold=True)
add_line(slide, Inches(0.8), Inches(1.65), Inches(1.5), ACCENT, 2)

add_text(slide, Inches(0.8), Inches(1.9), Inches(10), Inches(0.8),
         company["journey"]["collaborations"]["en"][:280] + "…",
         size=11, color=LIGHT_GRAY)

# Client cards — 4 columns
card_w = Inches(2.7)
card_h = Inches(3.2)
start_x = Inches(0.8)
start_y = Inches(3.2)
gap = Inches(0.3)

for i, client in enumerate(clients):
    x = start_x + i * (card_w + gap)
    # Card bg
    add_shape(slide, x, start_y, card_w, card_h, BG_CARD, DIVIDER)

    # Logo
    logo_path = safe_image(client["logo"])
    if logo_path:
        try:
            slide.shapes.add_picture(logo_path, x + Inches(0.6), start_y + Inches(0.35), Inches(1.5), Inches(1.0))
        except Exception:
            pass

    # Name
    add_text(slide, x + Inches(0.2), start_y + Inches(1.5), card_w - Inches(0.4), Inches(0.5),
             client["name"], size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Sector
    add_text(slide, x + Inches(0.2), start_y + Inches(2.0), card_w - Inches(0.4), Inches(0.4),
             client["sector"]["en"], size=10, color=ACCENT, alignment=PP_ALIGN.CENTER)

    # Divider in card
    add_line(slide, x + Inches(0.8), start_y + Inches(2.5), Inches(1.1), DIVIDER, 1)

    # Website
    site = client["website"].replace("https://", "").replace("http://", "").rstrip("/")
    add_text(slide, x + Inches(0.2), start_y + Inches(2.6), card_w - Inches(0.4), Inches(0.35),
             site, size=8, color=MUTED, alignment=PP_ALIGN.CENTER)

add_line(slide, Inches(0), H - Pt(3), W, ACCENT, 3)


# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — CLIENT EXPERTISE & COLLABORATION DETAIL
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, Inches(0), Inches(0), W, ACCENT, 3)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
         "CLIENT EXPERTISE", size=11, color=ACCENT, bold=True)
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
         "What We Deliver to Our Partners", size=32, color=WHITE, bold=True)
add_line(slide, Inches(0.8), Inches(1.65), Inches(1.5), ACCENT, 2)

# Specializations list (left column)
specs = company["journey"]["specializations"]
for i, spec in enumerate(specs):
    y = Inches(2.2) + i * Inches(0.7)
    add_accent_bar(slide, Inches(0.8), y, Inches(0.45))
    add_text(slide, Inches(1.1), y, Inches(5), Inches(0.45),
             spec["en"], size=13, color=LIGHT_GRAY)

# Right column — key stats
stats = company["stats"]
stat_x = Inches(7.5)
add_shape(slide, stat_x, Inches(2.2), Inches(4.8), Inches(4.2), BG_CARD, DIVIDER)

add_text(slide, stat_x + Inches(0.4), Inches(2.4), Inches(4), Inches(0.5),
         "AT A GLANCE", size=11, color=ACCENT, bold=True)

for i, stat in enumerate(stats):
    sy = Inches(3.0) + i * Inches(0.8)
    add_text(slide, stat_x + Inches(0.4), sy, Inches(1.5), Inches(0.5),
             stat["value"], size=28, color=ACCENT, bold=True, font_name="Segoe UI Light")
    add_text(slide, stat_x + Inches(2.2), sy + Inches(0.05), Inches(2.4), Inches(0.4),
             stat["label"]["en"], size=12, color=LIGHT_GRAY)

add_line(slide, Inches(0), H - Pt(3), W, ACCENT, 3)


# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — ABOUT / COMPANY STORY
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, Inches(0), Inches(0), W, ACCENT, 3)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
         "ABOUT ART SER", size=11, color=ACCENT, bold=True)
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
         "Our Story & Professional Journey", size=32, color=WHITE, bold=True)
add_line(slide, Inches(0.8), Inches(1.65), Inches(1.5), ACCENT, 2)

# Journey sections
sections = [
    ("THE BEGINNINGS", company["journey"]["beginnings"]["en"]),
    ("INDEPENDENCE", company["journey"]["independence"]["en"]),
    ("CURRENT WORK", company["journey"]["currentWork"]["en"][:220] + "…"),
]

for i, (heading, body) in enumerate(sections):
    y = Inches(2.1) + i * Inches(1.7)
    add_accent_bar(slide, Inches(0.8), y, Inches(1.2))
    add_text(slide, Inches(1.2), y, Inches(4), Inches(0.35),
             heading, size=10, color=ACCENT, bold=True)
    add_text(slide, Inches(1.2), y + Inches(0.35), Inches(10.5), Inches(0.9),
             body, size=11, color=LIGHT_GRAY)

add_line(slide, Inches(0), H - Pt(3), W, ACCENT, 3)


# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — SERVICES
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, Inches(0), Inches(0), W, ACCENT, 3)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
         "OUR SERVICES", size=11, color=ACCENT, bold=True)
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
         "Comprehensive Aluminium & Serramenti Solutions", size=28, color=WHITE, bold=True)
add_line(slide, Inches(0.8), Inches(1.65), Inches(1.5), ACCENT, 2)

# 3x2 grid
cols, rows = 3, 2
svc_w = Inches(3.7)
svc_h = Inches(2.3)
svc_gap_x = Inches(0.35)
svc_gap_y = Inches(0.35)
svc_start_x = Inches(0.8)
svc_start_y = Inches(2.1)

for idx, svc in enumerate(services):
    col = idx % cols
    row = idx // cols
    sx = svc_start_x + col * (svc_w + svc_gap_x)
    sy = svc_start_y + row * (svc_h + svc_gap_y)

    add_shape(slide, sx, sy, svc_w, svc_h, BG_CARD, DIVIDER)
    # Number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sx + Inches(0.2), sy + Inches(0.2), Inches(0.45), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()
    tf = badge.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = f"0{idx+1}"
    p.font.size = Pt(12)
    p.font.color.rgb = BG_DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text(slide, sx + Inches(0.8), sy + Inches(0.2), svc_w - Inches(1), Inches(0.45),
             svc["title"]["en"], size=12, color=WHITE, bold=True)
    add_text(slide, sx + Inches(0.2), sy + Inches(0.8), svc_w - Inches(0.4), Inches(1.3),
             svc["description"]["en"][:180] + ("…" if len(svc["description"]["en"]) > 180 else ""),
             size=9, color=MUTED)

add_line(slide, Inches(0), H - Pt(3), W, ACCENT, 3)


# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — FEATURED PROJECTS
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, Inches(0), Inches(0), W, ACCENT, 3)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
         "FEATURED PROJECTS", size=11, color=ACCENT, bold=True)
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
         "Precision Engineering in Action", size=32, color=WHITE, bold=True)
add_line(slide, Inches(0.8), Inches(1.65), Inches(1.5), ACCENT, 2)

for i, proj in enumerate(projects):
    px = Inches(0.8) + i * Inches(6.2)
    py = Inches(2.2)
    pw = Inches(5.6)
    ph = Inches(4.5)

    add_shape(slide, px, py, pw, ph, BG_CARD, DIVIDER)

    # Project image
    img = safe_image(proj["image"])
    if img:
        try:
            slide.shapes.add_picture(img, px + Inches(0.2), py + Inches(0.2), pw - Inches(0.4), Inches(2.2))
        except Exception:
            pass

    add_text(slide, px + Inches(0.3), py + Inches(2.5), pw - Inches(0.6), Inches(0.4),
             proj["title"]["en"], size=14, color=WHITE, bold=True)
    add_text(slide, px + Inches(0.3), py + Inches(2.95), pw - Inches(0.6), Inches(0.3),
             f'{proj["location"]}  ·  {proj["year"]}  ·  {proj["category"]["en"]}',
             size=9, color=ACCENT)
    add_text(slide, px + Inches(0.3), py + Inches(3.35), pw - Inches(0.6), Inches(1.0),
             proj["summary"]["en"][:250] + "…", size=9, color=MUTED)

add_line(slide, Inches(0), H - Pt(3), W, ACCENT, 3)


# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — SUPPLIERS / PARTNERS
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, Inches(0), Inches(0), W, ACCENT, 3)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
         "OUR SUPPLIERS", size=11, color=ACCENT, bold=True)
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
         "Premium European Material Partners", size=32, color=WHITE, bold=True)
add_line(slide, Inches(0.8), Inches(1.65), Inches(1.5), ACCENT, 2)

# Supplier cards — flexible row
sup_w = Inches(2.2)
sup_h = Inches(4.0)
sup_start_x = Inches(0.5)
sup_start_y = Inches(2.2)
sup_gap = Inches(0.25)

for i, sup in enumerate(suppliers):
    sx = sup_start_x + i * (sup_w + sup_gap)

    add_shape(slide, sx, sup_start_y, sup_w, sup_h, BG_CARD, DIVIDER)

    # Logo
    logo = safe_image(sup["logo"])
    if logo:
        try:
            slide.shapes.add_picture(logo, sx + Inches(0.35), sup_start_y + Inches(0.25), Inches(1.5), Inches(0.9))
        except Exception:
            pass

    add_text(slide, sx + Inches(0.15), sup_start_y + Inches(1.3), sup_w - Inches(0.3), Inches(0.35),
             sup["name"], size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, sx + Inches(0.15), sup_start_y + Inches(1.65), sup_w - Inches(0.3), Inches(0.25),
             f'{sup["country"]}  ·  {sup["type"]["en"]}', size=7, color=ACCENT, alignment=PP_ALIGN.CENTER)

    # Products
    prods = ", ".join(sup["products"][:4])
    add_text(slide, sx + Inches(0.15), sup_start_y + Inches(2.1), sup_w - Inches(0.3), Inches(1.5),
             prods, size=7, color=MUTED, alignment=PP_ALIGN.CENTER)

add_line(slide, Inches(0), H - Pt(3), W, ACCENT, 3)


# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — CERTIFICATIONS
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, Inches(0), Inches(0), W, ACCENT, 3)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
         "CERTIFICATIONS & QUALIFICATIONS", size=11, color=ACCENT, bold=True)
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
         "Professional Standards & Safety Compliance", size=28, color=WHITE, bold=True)
add_line(slide, Inches(0.8), Inches(1.65), Inches(1.5), ACCENT, 2)

# 3x2 grid
cert_w = Inches(3.7)
cert_h = Inches(2.2)

for idx, cert in enumerate(certs):
    col = idx % 3
    row = idx // 3
    cx = Inches(0.8) + col * (cert_w + Inches(0.35))
    cy = Inches(2.1) + row * (cert_h + Inches(0.35))

    add_shape(slide, cx, cy, cert_w, cert_h, BG_CARD, DIVIDER)

    # Year badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx + Inches(0.2), cy + Inches(0.2), Inches(0.6), Inches(0.35))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = str(cert["year"])
    p.font.size = Pt(10)
    p.font.color.rgb = BG_DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text(slide, cx + Inches(1.0), cy + Inches(0.2), cert_w - Inches(1.2), Inches(0.35),
             cert["name"], size=11, color=WHITE, bold=True)
    add_text(slide, cx + Inches(0.2), cy + Inches(0.7), cert_w - Inches(0.4), Inches(0.3),
             cert["issuer"], size=9, color=ACCENT)
    add_text(slide, cx + Inches(0.2), cy + Inches(1.1), cert_w - Inches(0.4), Inches(0.9),
             cert["description"]["en"], size=9, color=MUTED)

add_line(slide, Inches(0), H - Pt(3), W, ACCENT, 3)


# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — MISSION & VISION
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, Inches(0), Inches(0), W, ACCENT, 3)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
         "MISSION & VISION", size=11, color=ACCENT, bold=True)
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
         "Our Guiding Principles", size=32, color=WHITE, bold=True)
add_line(slide, Inches(0.8), Inches(1.65), Inches(1.5), ACCENT, 2)

# Mission card
mx = Inches(0.8)
my = Inches(2.3)
mw = Inches(5.5)
mh = Inches(3.5)
add_shape(slide, mx, my, mw, mh, BG_CARD, DIVIDER)
add_accent_bar(slide, mx + Inches(0.3), my + Inches(0.3), Inches(0.5))
add_text(slide, mx + Inches(0.6), my + Inches(0.3), Inches(3), Inches(0.5),
         "OUR MISSION", size=14, color=ACCENT, bold=True)
add_text(slide, mx + Inches(0.3), my + Inches(1.0), mw - Inches(0.6), Inches(2.2),
         company["mission"]["en"], size=13, color=LIGHT_GRAY)

# Vision card
vx = Inches(7.0)
add_shape(slide, vx, my, mw, mh, BG_CARD, DIVIDER)
add_accent_bar(slide, vx + Inches(0.3), my + Inches(0.3), Inches(0.5))
add_text(slide, vx + Inches(0.6), my + Inches(0.3), Inches(3), Inches(0.5),
         "OUR VISION", size=14, color=ACCENT, bold=True)
add_text(slide, vx + Inches(0.3), my + Inches(1.0), mw - Inches(0.6), Inches(2.2),
         company["vision"]["en"], size=13, color=LIGHT_GRAY)

add_line(slide, Inches(0), H - Pt(3), W, ACCENT, 3)


# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — CONTACT / CLOSING
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, Inches(0), Inches(0), W, ACCENT, 4)

add_text(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1),
         "ART SER", size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Segoe UI Light")

add_line(slide, Inches(5.5), Inches(3.2), Inches(2.3), ACCENT, 2)

add_text(slide, Inches(2), Inches(3.6), Inches(9), Inches(0.5),
         "Let's Build Something Exceptional Together", size=20, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(2), Inches(4.5), Inches(9), Inches(0.4),
         "ART SER DI SHEHEZAD TARIQ", size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER, bold=True)

add_text(slide, Inches(2), Inches(5.0), Inches(9), Inches(0.4),
         "Verona, Italy  ·  26+ Years of Industry Excellence", size=12, color=MUTED, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(2), Inches(5.8), Inches(9), Inches(0.4),
         "www.artser.it", size=11, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_line(slide, Inches(0), H - Pt(4), W, ACCENT, 4)


# ── Save ──
prs.save(OUT)
print(f"Presentation saved -> {OUT}")
print(f"  Slides: {len(prs.slides)}")
